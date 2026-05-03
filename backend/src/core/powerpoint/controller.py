"""
PowerPoint Controller Module

This module provides the main controller for PowerPoint automation operations,
including presentation creation from images and video export functionality.
"""

import os
import re
from pathlib import Path
from typing import List, Optional, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt
from PIL import Image


# Error Classes
class PowerPointError(Exception):
    """Base exception for PowerPoint operations."""
    pass


class TemplateError(PowerPointError):
    """Template file is invalid or corrupted."""
    pass


class PowerPointNotFoundError(PowerPointError):
    """PowerPoint is not installed or not accessible."""
    pass


class ExportError(PowerPointError):
    """Video export operation failed."""
    pass


class OperationCancelledError(PowerPointError):
    """Operation was cancelled by user."""
    pass


class PowerPointController:
    """Main controller for PowerPoint automation operations."""
    
    def __init__(self, config: Optional[Dict[str, Any]] = None):
        """
        Initialize with configuration settings.
        
        Args:
            config: Dictionary containing PowerPoint configuration settings.
                   If None, uses default configuration.
        """
        self.config = config or self._get_default_config()
        
    def _get_default_config(self) -> Dict[str, Any]:
        """
        Get default configuration settings.
        
        Returns:
            Dictionary with default PowerPoint settings
        """
        return {
            'template_path': 'templates/default.pptm',
            'output_folder': 'output/presentations',
            'video_folder': 'output/videos',
            'slide_duration': 3.0,
            'transition_type': 'fade',
            'transition_duration': 0.5,
            'video_resolution': (3840, 2160),
            'video_fps': 30,
            'video_quality': 5,
            'image_fit_mode': 'contain',
            'image_position': 'center',
            'preserve_aspect_ratio': True
        }
    
    def get_image_files(self, folder_path: str) -> List[str]:
        """
        Get sorted list of image files from folder.
        
        Detects image files with common extensions (png, jpg, jpeg, gif, bmp)
        and sorts them by numeric prefix if present, otherwise alphabetically.
        
        Args:
            folder_path: Path to folder containing images
            
        Returns:
            List of image file paths sorted appropriately
            
        Raises:
            FileNotFoundError: If folder doesn't exist
            ValueError: If folder is empty or contains no images
        """
        folder = Path(folder_path)
        
        if not folder.exists():
            raise FileNotFoundError(f"Image folder not found: {folder_path}")
        
        if not folder.is_dir():
            raise ValueError(f"Path is not a directory: {folder_path}")
        
        # Supported image extensions
        image_extensions = {'.png', '.jpg', '.jpeg', '.gif', '.bmp'}
        
        # Get all image files
        image_files = [
            str(f) for f in folder.iterdir()
            if f.is_file() and f.suffix.lower() in image_extensions
        ]
        
        if not image_files:
            raise ValueError(f"No image files found in folder: {folder_path}")
        
        # Sort by numeric prefix if present, otherwise alphabetically
        def sort_key(filepath: str) -> tuple:
            """
            Generate sort key for image file.
            
            Extracts numeric prefix from filename for sorting.
            Files with numeric prefixes sort before those without.
            """
            filename = Path(filepath).name
            # Try to extract leading number from filename
            match = re.match(r'^(\d+)', filename)
            if match:
                # Return tuple: (has_number=0, number, filename)
                # has_number=0 ensures numbered files come first
                return (0, int(match.group(1)), filename)
            else:
                # Return tuple: (has_number=1, 0, filename)
                # has_number=1 ensures non-numbered files come last
                return (1, 0, filename)
        
        image_files.sort(key=sort_key)
        
        return image_files
    
    def validate_template(self, template_path: str) -> bool:
        """
        Validate that template file is usable.
        
        Checks if the template file exists, has the correct extension,
        and is a valid PowerPoint file.
        
        Args:
            template_path: Path to .pptx template file
            
        Returns:
            True if template is valid
            
        Raises:
            FileNotFoundError: If template file doesn't exist
            ValueError: If template file is not a valid .pptx file
        """
        template = Path(template_path)
        
        # Check if file exists
        if not template.exists():
            raise FileNotFoundError(f"Template file not found: {template_path}")
        
        # Check if it's a file (not a directory)
        if not template.is_file():
            raise ValueError(f"Template path is not a file: {template_path}")
        
        # Check file extension
        if template.suffix.lower() not in ('.pptx', '.pptm'):
            raise ValueError(
                f"Invalid template format. Expected .pptx or .pptm file, got: {template.suffix}"
            )
        
        # Check file size (basic validation - should be at least 10KB for a valid pptx)
        file_size = template.stat().st_size
        if file_size < 10240:  # 10KB minimum
            raise ValueError(
                f"Template file appears to be corrupted or invalid (too small: {file_size} bytes)"
            )
        
        # Basic PPTX format validation (check for ZIP signature)
        # PPTX files are ZIP archives, so they should start with PK signature
        try:
            with open(template_path, 'rb') as f:
                header = f.read(4)
                if header[:2] != b'PK':
                    raise ValueError(
                        "Template file is not a valid PowerPoint file (invalid ZIP signature)"
                    )
        except IOError as e:
            raise ValueError(f"Cannot read template file: {e}")
        
        return True
    
    def create_presentation(
        self,
        template_path: str,
        output_path: str,
        image_files: Optional[List[str]] = None,
        image_folder: Optional[str] = None,
        slide_duration: float = 3.0,
        transition_type: str = "fade"
    ) -> str:
        """
        Create a PowerPoint presentation from images.

        Pass either an explicit list of image paths via ``image_files`` (the
        preferred path — guarantees the deck only contains screenshots from
        the *current* run) or a folder via ``image_folder`` (legacy behaviour
        that globs every PNG in the folder, including leftovers from previous
        runs). Mixing the two raises ``ValueError``.

        Args:
            template_path: Path to .pptx template file
            output_path: Where to save the presentation
            image_files: Explicit ordered list of image paths to insert.
            image_folder: Folder to glob (legacy fallback only).
            slide_duration: Duration per slide in seconds
            transition_type: Type of transition between slides

        Returns:
            Path to created presentation file

        Raises:
            FileNotFoundError: If template or images not found
            ValueError: If invalid configuration provided
            PowerPointError: If PowerPoint operation fails
        """
        # Validate template
        self.validate_template(template_path)

        # Resolve image list — prefer explicit ``image_files`` so a deck
        # only contains screenshots from this run, never leftovers in the
        # output folder from a previous job. Fall back to globbing only
        # when the caller explicitly opts in via ``image_folder``.
        if image_files is not None and image_folder is not None:
            raise ValueError("Pass either image_files or image_folder, not both")
        if image_files is not None:
            if not image_files:
                raise ValueError("image_files is empty")
            missing = [p for p in image_files if not Path(p).is_file()]
            if missing:
                raise FileNotFoundError(
                    f"Image file(s) not found: {missing[0]}"
                    + (f" (+{len(missing) - 1} more)" if len(missing) > 1 else "")
                )
            resolved_image_files: List[str] = list(image_files)
        elif image_folder is not None:
            resolved_image_files = self.get_image_files(image_folder)
        else:
            raise ValueError("Must provide either image_files or image_folder")
        image_files = resolved_image_files
        
        # Validate slide duration
        if slide_duration <= 0:
            raise ValueError(f"Slide duration must be positive, got: {slide_duration}")
        
        # Load template
        try:
            prs = Presentation(template_path)
        except Exception as e:
            raise TemplateError(f"Failed to load template: {e}")
        
        # Get slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height
        
        # Get blank slide layout (typically the last layout or layout index 6)
        # Try to find a blank layout, otherwise use the first available layout
        try:
            blank_layout = prs.slide_layouts[6]  # Blank layout is usually at index 6
        except IndexError:
            # If index 6 doesn't exist, use the last layout
            blank_layout = prs.slide_layouts[-1]
        
        # Create slides for each image
        for image_path in image_files:
            try:
                # Add a new slide with blank layout
                slide = prs.slides.add_slide(blank_layout)
                
                # Insert image centered with aspect ratio preservation
                self._insert_image_centered(
                    slide, 
                    image_path, 
                    slide_width, 
                    slide_height
                )
                
                # Set slide duration and transition
                self._set_slide_timing_and_transition(slide, slide_duration, transition_type)
                
            except Exception as e:
                raise PowerPointError(f"Failed to process image {image_path}: {e}")
        
        # Ensure output directory exists
        output_dir = Path(output_path).parent
        output_dir.mkdir(parents=True, exist_ok=True)
        
        # Save presentation
        try:
            prs.save(output_path)
        except Exception as e:
            raise PowerPointError(f"Failed to save presentation: {e}")
        
        return output_path
    
    def _insert_image_centered(
        self, 
        slide, 
        image_path: str, 
        slide_width: int, 
        slide_height: int
    ) -> None:
        """
        Insert image centered on slide, preserving aspect ratio.
        
        Args:
            slide: PowerPoint slide object
            image_path: Path to image file
            slide_width: Width of slide in EMUs
            slide_height: Height of slide in EMUs
        """
        # Open image to get dimensions
        try:
            img = Image.open(image_path)
            img_width, img_height = img.size
        except Exception as e:
            raise PowerPointError(f"Failed to open image {image_path}: {e}")
        
        # Calculate aspect ratios
        img_aspect = img_width / img_height
        slide_aspect = slide_width / slide_height
        
        # Calculate dimensions to fit slide while preserving aspect ratio
        if img_aspect > slide_aspect:
            # Image is wider - fit to width
            width = slide_width
            height = int(width / img_aspect)
        else:
            # Image is taller - fit to height
            height = slide_height
            width = int(height * img_aspect)
        
        # Center the image
        left = (slide_width - width) // 2
        top = (slide_height - height) // 2
        
        # Add picture to slide
        try:
            slide.shapes.add_picture(
                image_path,
                left, top,
                width, height
            )
        except Exception as e:
            raise PowerPointError(f"Failed to add image to slide: {e}")
    
    def _set_slide_timing_and_transition(
        self, 
        slide, 
        slide_duration: float,
        transition_type: str
    ) -> None:
        """
        Set slide timing and transition type.
        
        Args:
            slide: PowerPoint slide object
            slide_duration: Duration in seconds before auto-advance
            transition_type: Type of transition (fade, push, wipe, none, etc.)
        """
        from lxml import etree
        
        # Get the slide element
        slide_elem = slide._element
        
        # Define namespaces
        nsmap = {
            'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
            'a': 'http://schemas.openxmlformats.org/drawingml/2006/main'
        }
        
        # Check if transition element exists, if not create it
        transition = slide_elem.find('p:transition', nsmap)
        if transition is None:
            transition = etree.SubElement(slide_elem, '{%s}transition' % nsmap['p'])
        
        # Set auto-advance timing (in milliseconds)
        # advTm attribute controls automatic advance after specified time
        transition.set('advTm', str(int(slide_duration * 1000)))
        
        # Set to advance automatically (not on click)
        transition.set('advClick', '0')  # 0 = don't advance on click, 1 = advance on click
        
        # Apply transition type
        # Remove any existing transition effect elements
        for child in list(transition):
            transition.remove(child)
        
        # Map transition types to PowerPoint XML elements
        transition_type_lower = transition_type.lower()
        
        if transition_type_lower == 'fade':
            # Add fade transition
            fade_elem = etree.SubElement(transition, '{%s}fade' % nsmap['p'])
            fade_elem.set('thruBlk', '0')  # Fade through black: 0 = no, 1 = yes
        elif transition_type_lower == 'push':
            # Add push transition
            push_elem = etree.SubElement(transition, '{%s}push' % nsmap['p'])
            push_elem.set('dir', 'l')  # Direction: l = left, r = right, u = up, d = down
        elif transition_type_lower == 'wipe':
            # Add wipe transition
            wipe_elem = etree.SubElement(transition, '{%s}wipe' % nsmap['p'])
            wipe_elem.set('dir', 'l')  # Direction: l = left, r = right, u = up, d = down
        elif transition_type_lower == 'none' or transition_type_lower == '':
            # No transition effect - just timing
            pass
        else:
            # Default to fade for unknown transition types
            fade_elem = etree.SubElement(transition, '{%s}fade' % nsmap['p'])
            fade_elem.set('thruBlk', '0')
    def export_to_video(
        self,
        presentation_path: str,
        output_video_path: str,
        resolution: tuple = (3840, 2160),
        fps: int = 30,
        quality: int = 5,
        progress_callback=None,
        cancel_event=None,
    ) -> str:
        """
        Export presentation to 4K video using COM automation.

        Args:
            presentation_path: Path to .pptx file
            output_video_path: Where to save video file
            resolution: Video resolution (width, height)
            fps: Frames per second
            quality: Video quality (1-5, 5 is highest)

        Returns:
            Path to exported video file

        Raises:
            FileNotFoundError: If presentation file doesn't exist
            PowerPointNotFoundError: If PowerPoint not installed
            ExportError: If video export fails
        """
        # Import exporter here to avoid import errors on non-Windows systems
        try:
            from src.core.powerpoint.exporter import PowerPointExporter
        except ImportError as e:
            raise PowerPointNotFoundError(
                "PowerPoint exporter not available. "
                "This feature requires Windows and pywin32."
            ) from e

        # Validate presentation file exists
        if not Path(presentation_path).exists():
            raise FileNotFoundError(
                f"Presentation file not found: {presentation_path}"
            )

        # Validate quality parameter
        if not 1 <= quality <= 5:
            raise ValueError(
                f"Video quality must be between 1 and 5, got: {quality}"
            )

        # Validate resolution
        if len(resolution) != 2 or resolution[0] <= 0 or resolution[1] <= 0:
            raise ValueError(
                f"Invalid resolution. Expected (width, height) with positive values, "
                f"got: {resolution}"
            )

        # Validate fps
        if fps <= 0:
            raise ValueError(f"FPS must be positive, got: {fps}")

        # Extract width and height from resolution tuple
        width, height = resolution

        # Create exporter and export video with proper cleanup
        exporter = None
        try:
            # Create PowerPointExporter instance
            exporter = PowerPointExporter()

            # Check if PowerPoint is installed
            if not exporter.is_powerpoint_installed():
                raise PowerPointNotFoundError(
                    "Microsoft PowerPoint is not installed or not accessible. "
                    "Video export requires PowerPoint to be installed on Windows."
                )

            # Open presentation via COM
            exporter.open_presentation(presentation_path)

            # Export to MP4 with specified settings
            exporter.export_video(
                output_path=output_video_path,
                width=width,
                height=height,
                fps=fps,
                quality=quality,
                progress_callback=progress_callback,
                cancel_event=cancel_event,
            )

            return output_video_path

        except FileNotFoundError:
            # Re-raise file not found errors as-is
            raise
        except ValueError:
            # Re-raise validation errors as-is
            raise
        except PowerPointNotFoundError:
            # Re-raise PowerPoint not found errors as-is
            raise
        except Exception as e:
            # Wrap other exceptions as ExportError
            raise ExportError(
                f"Failed to export video: {e}"
            ) from e
        finally:
            # Handle cleanup in try-finally block
            if exporter is not None:
                try:
                    exporter.quit_powerpoint()
                except Exception as cleanup_error:
                    # Log cleanup errors but don't raise them
                    # (we don't want cleanup errors to mask the original error)
                    import logging
                    logger = logging.getLogger(__name__)
                    logger.warning(
                        f"Error during PowerPoint cleanup: {cleanup_error}"
                    )

    def create_template_presentation(
        self,
        template_path: str,
        image_files: list,
        output_pptx_path: str,
        slide_duration: float = 3.0,
        base_slide_index: int = 3,
        intro_thumbnail_path: str | None = None,
        intro_thumbnail_duration: float = 5.0,
        outro_thumbnail_path: str | None = None,
        outro_thumbnail_duration: float = 5.0,
        progress_callback=None,
        cancel_event=None,
    ) -> str:
        """Create a deck through PowerPoint COM while preserving template slides."""
        try:
            from src.core.powerpoint.exporter import PowerPointExporter
        except ImportError as e:
            raise PowerPointNotFoundError(
                "PowerPoint exporter not available. "
                "This feature requires Windows and pywin32."
            ) from e

        exporter = None
        try:
            exporter = PowerPointExporter()
            if not exporter.is_powerpoint_installed():
                raise PowerPointNotFoundError(
                    "Microsoft PowerPoint is not installed or not accessible."
                )
            return exporter.create_from_template(
                template_path=template_path,
                image_files=image_files,
                output_path=output_pptx_path,
                base_slide_index=base_slide_index,
                slide_duration=slide_duration,
                intro_thumbnail_path=intro_thumbnail_path,
                intro_thumbnail_duration=intro_thumbnail_duration,
                outro_thumbnail_path=outro_thumbnail_path,
                outro_thumbnail_duration=outro_thumbnail_duration,
                progress_callback=progress_callback,
                cancel_event=cancel_event,
            )
        except PowerPointNotFoundError:
            raise
        except Exception as e:
            raise ExportError(f"Failed to create presentation: {e}") from e
        finally:
            if exporter is not None:
                try:
                    exporter.quit_powerpoint()
                except Exception as cleanup_error:
                    import logging
                    logger = logging.getLogger(__name__)
                    logger.warning(
                        f"Error during PowerPoint cleanup: {cleanup_error}"
                    )

    def create_and_export_video(
        self,
        template_path: str,
        image_files: list,
        output_pptx_path: str,
        output_video_path: str,
        slide_duration: float = 3.0,
        transition_type: str = "fade",
        resolution: tuple = (3840, 2160),
        fps: int = 30,
        quality: int = 5,
        base_slide_index: int = 3,
        intro_thumbnail_path: str | None = None,
        intro_thumbnail_duration: float = 5.0,
        outro_thumbnail_path: str | None = None,
        outro_thumbnail_duration: float = 5.0,
        progress_callback=None,
        cancel_event=None
    ) -> dict:
        """
        Full workflow: Create presentation from template + export to video.
        
        Uses COM automation for the entire pipeline in one session,
        preserving template watermarks, transitions, and effects.
        
        Args:
            template_path: Path to .pptm template
            image_files: List of screenshot file paths 
            output_pptx_path: Where to save the presentation
            output_video_path: Where to save the video
            slide_duration: Seconds per slide
            transition_type: Transition type (used only as fallback)
            resolution: Video resolution (width, height)
            fps: Frames per second
            quality: Video quality 1-5
            base_slide_index: 1-based index of content base slide (default: 3)
            
        Returns:
            Dict with 'presentation_path' and 'video_path' (None if export failed)
        """
        try:
            from src.core.powerpoint.exporter import PowerPointExporter
        except ImportError as e:
            raise PowerPointNotFoundError(
                "PowerPoint exporter not available. "
                "This feature requires Windows and pywin32."
            ) from e
        
        result = {
            'presentation_path': None,
            'video_path': None,
            'warning': None
        }
        
        exporter = None
        try:
            exporter = PowerPointExporter()
            
            if not exporter.is_powerpoint_installed():
                raise PowerPointNotFoundError(
                    "Microsoft PowerPoint is not installed or not accessible."
                )
            
            # Step 1: Create presentation from template via COM
            exporter.create_from_template(
                template_path=template_path,
                image_files=image_files,
                output_path=output_pptx_path,
                base_slide_index=base_slide_index,
                slide_duration=slide_duration,
                intro_thumbnail_path=intro_thumbnail_path,
                intro_thumbnail_duration=intro_thumbnail_duration,
                outro_thumbnail_path=outro_thumbnail_path,
                outro_thumbnail_duration=outro_thumbnail_duration,
                progress_callback=progress_callback,
                cancel_event=cancel_event
            )
            result['presentation_path'] = output_pptx_path
            
            # Step 2: Export to video (presentation is already open in COM)
            width, height = resolution
            
            # Ensure video output directory exists
            from pathlib import Path as PPath
            PPath(output_video_path).parent.mkdir(parents=True, exist_ok=True)
            
            exporter.export_video(
                output_path=output_video_path,
                width=width,
                height=height,
                fps=fps,
                quality=quality,
                default_slide_duration=slide_duration,
                progress_callback=progress_callback,
                cancel_event=cancel_event
            )
            if Path(output_video_path).exists() and Path(output_video_path).stat().st_size > 0:
                result['video_path'] = output_video_path
            else:
                result['warning'] = "PowerPoint reported completion, but the MP4 file was not found."
            
            return result
            
        except PowerPointNotFoundError:
            raise
        except Exception as e:
            # If we at least created the presentation, return that
            if result['presentation_path'] and Path(output_pptx_path).exists():
                result['warning'] = f"Video export failed: {str(e)}"
                return result
            raise ExportError(f"Failed to create presentation and video: {e}") from e
        finally:
            if exporter is not None:
                try:
                    exporter.quit_powerpoint()
                except Exception as cleanup_error:
                    import logging
                    logger = logging.getLogger(__name__)
                    logger.warning(
                        f"Error during PowerPoint cleanup: {cleanup_error}"
                    )
